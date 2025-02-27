import streamlit as st
import fitz  # PyMuPDF for PDF text extraction
from groq import Groq
import json
import requests
from streamlit_lottie import st_lottie
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from dotenv import load_dotenv
import os

GROQ_API_KEY = "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"

client = Groq(api_key=GROQ_API_KEY)

st.set_page_config(page_title="📜 Research Paper Summarizer", layout="wide")

def load_lottieurl(url):
    r = requests.get(url)
    if r.status_code != 200:
        return None
    return r.json()

upload_animation = load_lottieurl("https://assets2.lottiefiles.com/packages/lf20_s4tubmwg.json")

dark_theme = """
    <style>
        body {
            background-color: #1E1E1E;
            color: #FFFFFF;
        }
        .stTextInput > div > div > input {
            background-color: #2D2D2D;
            color: #FFFFFF;
        }
        .stFileUploader > div {
            background-color: #2D2D2D;
            color: #FFFFFF;
        }
        .stButton>button {
            background-color: #ff6347;
            color: white;
            border-radius: 10px;
            padding: 10px 20px;
            font-size: 18px;
            border: none;
        }
        .stButton>button:hover {
            background-color: #ff4500;
        }
    </style>
"""
st.markdown(dark_theme, unsafe_allow_html=True)

st.markdown(
    "<h1 style='text-align: center; color: #ff6347;'>📜 Research Paper Summarizer 🤖</h1>",
    unsafe_allow_html=True,
)
st.markdown(
    "<p style='text-align: center; font-size:18px;'>Upload a PDF research paper, and I'll generate a short summary for you! 🚀</p>",
    unsafe_allow_html=True,
)

st_lottie(upload_animation, height=150, key="upload")

uploaded_file = st.file_uploader("📂 Upload a Research Paper (PDF)", type="pdf")

def extract_text_from_pdf(file):
    text = ""
    with fitz.open(stream=file.read()) as doc:
        for page in doc:
            text += page.get_text("text") + "\n"
    return text

def summarize_with_groq(text):
    try:
        chat_completion = client.chat.completions.create(
            messages=[
                {"role": "system", "content": "Summarize this research paper in a structured format with key points."},
                {"role": "user", "content": text},
            ],
            model="deepseek-r1-distill-llama-70b",
        )
        return chat_completion.choices[0].message.content
    except Exception as e:
        return f"❌ Error: {str(e)}"

def generate_pptx(lesson_topic):
    try:
        prompt = (
            "Create PowerPoint slides for a lesson plan based on this research paper summary. "
            "For each slide, provide a clear title and up to 4 concise bullet points. "
            "Do not include labels like 'Heading:', 'Sub-point:', etc. - just write the content directly. "
            "Start with a title slide and end with key takeaways. "
            "Format each slide as:\n\n"
            "[Title of Slide]\n"
            "• First point\n"
            "• Second point\n"
            "• Third point\n"
            "• Fourth point\n\n"
            "Research Paper Summary:\n{}"
        ).format(lesson_topic)

        chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": "You are a helpful assistant that creates clear and concise PowerPoint slide outlines for teaching purposes. Present content directly without labels or prefixes."
                },
                {"role": "user", "content": prompt},
            ],
            model="llama-3.3-70b-versatile",
        )
        return chat_completion.choices[0].message.content
    except Exception as e:
        return f"❌ Error generating PowerPoint outline: {str(e)}"

def apply_slide_styling(slide):
    # Set background color (light gray for professional look)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Style the title if it exists
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_frame = title_shape.text_frame
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 44, 44)

def create_ppt(summary_text):
    lesson_topic = summary_text.strip()
    ppt_outline = generate_pptx(lesson_topic)
    
    # Create a PowerPoint presentation based on the outline
    prs = Presentation()
    slides = ppt_outline.strip().split('\n\n')
    
    for slide_content in slides:
        lines = [line.strip() for line in slide_content.split('\n') if line.strip()]
        if not lines:
            continue
            
        # First line is the title
        slide_title = lines[0].strip('[]')  # Remove brackets if present
        bullet_points = lines[1:]  # Rest are bullet points

        # Add a slide to the presentation
        slide_layout = prs.slide_layouts[1]
        current_slide = prs.slides.add_slide(slide_layout)
        
        # Apply styling to the slide
        apply_slide_styling(current_slide)
        
        # Set title
        title = current_slide.shapes.title
        title.text = slide_title
        
        # Add and style content
        if bullet_points:
            content = current_slide.placeholders[1]
            text_frame = content.text_frame
            
            # Add each bullet point
            for idx, point in enumerate(bullet_points):
                p = text_frame.add_paragraph()
                p.text = point.strip('• ')  # Remove bullet if present
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(66, 66, 66)
                p.space_after = Pt(12)
                p.space_before = Pt(6)
                p.level = 0  # This ensures it's a bullet point

    # Save PPT to memory buffer
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

if uploaded_file:
    st.success("✅ File uploaded successfully!")
    pdf_text = extract_text_from_pdf(uploaded_file)

    if not pdf_text.strip():
        st.error("❌ No text found in the PDF! Please upload another file.")
        st.stop()
    
    if st.button("✨ Generate Summary & PPT", use_container_width=True):
        with st.spinner("⏳ Summarizing... please wait!"):
            summary = summarize_with_groq(pdf_text)
        
        st.subheader("🔍 Research Paper Summary")
        st.write(summary)
        
        ppt_file = create_ppt(summary)
        st.download_button(
            label="📥 Download Summary PPT",
            data=ppt_file,
            file_name="summary_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )
